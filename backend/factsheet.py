"""
Factsheet generator — python-pptx, A4 portrait (8.27" × 11.69").
Layout modelled on the Credicorp Capital template.
"""

import matplotlib
matplotlib.use("Agg")
matplotlib.rcParams['font.family'] = 'Inter'
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import pandas as pd
import numpy as np
from io import BytesIO
from datetime import datetime, date, timedelta

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    _PPTX = True
except ImportError:
    _PPTX = False

try:
    import yfinance as yf
    YF_OK = True
except ImportError:
    YF_OK = False

from backend.market_data import resolve_ticker

# ── Page geometry ──────────────────────────────────────────────────────────────
_W, _H   = 8.27, 11.69   # A4 portrait inches
_ML, _MR = 0.30, 0.30
_CW      = _W - _ML - _MR  # 7.67"

# Two-column split (main + bottom sections)
_LW  = 3.45                   # left col width
_GAP = 0.15
_RW  = _CW - _LW - _GAP       # right col width  ≈ 4.07"
_RX  = _ML + _LW + _GAP       # right col left-edge ≈ 3.90"

# Y positions: (top_y, height) in inches from slide top
_YP = {
    "header":    (0.00, 0.55),
    "title":     (0.60, 0.30),
    "det_bar":   (0.95, 0.22),
    "narrative": (1.21, 0.92),
    "sbars":     (2.18, 0.22),
    "main":      (2.44, 3.40),
    "sumtbl":    (5.89, 0.85),
    "evol_bar":  (6.79, 0.22),
    "subhdr":    (7.05, 0.30),
    "bottom":    (7.40, 3.00),
    "footnote":  (10.49, 0.18),
}

# ── Colour helpers ─────────────────────────────────────────────────────────────
_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT  = RGBColor(0xF2, 0xF4, 0xF8)
_DARK   = RGBColor(0x11, 0x18, 0x27)
_MUTED  = RGBColor(0x6B, 0x72, 0x80)

_LINE_COLORS = ["#1A1A2E", "#DC2626", "#9CA3AF", "#2563EB"]

_MESES_SHORT = ["Ene","Feb","Mar","Abr","May","Jun",
                "Jul","Ago","Set","Oct","Nov","Dic"]
_MESES_LONG  = ["enero","febrero","marzo","abril","mayo","junio",
                "julio","agosto","septiembre","octubre","noviembre","diciembre"]
_ORDINAL_ES  = ["primera","segunda","tercera","cuarta","quinta",
                "sexta","séptima","octava","novena","décima"]

_UND_LABELS = {
    "IWM": "Small Caps EE.UU. (IWM)",
    "XLF": "Sector Financiero (XLF)",
    "XLV": "Sector Salud (XLV)",
    "XLP": "Sector Consumo No Discrecional (XLP)",
    "SPX": "S&P 500 (SPX)", "NDX": "NASDAQ 100 (NDX)",
    "RTY": "Russell 2000 (RTY)", "SX5E": "Euro Stoxx 50 (SX5E)",
}


def _rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


# ── Date / number helpers ──────────────────────────────────────────────────────

_ES_MON = {"Ene":"Jan","Feb":"Feb","Mar":"Mar","Abr":"Apr","May":"May",
           "Jun":"Jun","Jul":"Jul","Ago":"Aug","Set":"Sep","Sep":"Sep",
           "Oct":"Oct","Nov":"Nov","Dic":"Dec"}

def _parse_date(s) -> date | None:
    if s is None:
        return None
    if isinstance(s, datetime):
        return s.date()
    if isinstance(s, date):
        return s
    norm = str(s).strip()
    if norm in ("", "nan", "None", "NaT"):
        return None
    for es, en in _ES_MON.items():
        norm = norm.replace(es, en)
    for fmt in ("%Y-%m-%d", "%d/%m/%Y", "%d-%b-%y", "%d-%b-%Y",
                "%d-%m-%Y", "%m/%d/%Y"):
        try:
            return datetime.strptime(norm, fmt).date()
        except ValueError:
            pass
    return None


def _fmt_short(d) -> str:
    d = _parse_date(d)
    return f"{d.day:02d}-{_MESES_SHORT[d.month-1]}-{str(d.year)[2:]}" if d else "—"


def _fmt_long(d) -> str:
    d = _parse_date(d)
    return f"{d.day} de {_MESES_LONG[d.month-1]} de {d.year}" if d else "—"


def _sf(v, default=None):
    try:
        f = float(v)
        return default if (f != f) else f
    except (TypeError, ValueError):
        return default


def _pct(v) -> str:
    f = _sf(v)
    if f is None:
        return "—"
    if abs(f) <= 1.5:
        f *= 100
    return f"{f:.2f}%"


# ── Slide drawing primitives ───────────────────────────────────────────────────

def _rect(slide, x, y, w, h, fill: RGBColor, border: RGBColor | None = None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.5)
    return shape


def _txt(slide, x, y, w, h, text, size=9, bold=False, color=None,
         align=PP_ALIGN.LEFT, italic=False):
    if color is None:
        color = _DARK
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name   = "Inter"
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _bar(slide, x, y, w, h, text, fill: RGBColor, size=9):
    shape = _rect(slide, x, y, w, h, fill)
    tf = shape.text_frame
    tf.margin_left   = Inches(0.06)
    tf.margin_right  = Inches(0.04)
    tf.margin_top    = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name  = "Inter"
    run.font.size  = Pt(size)
    run.font.bold  = True
    run.font.color.rgb = _WHITE


def _multiline_txt(slide, x, y, w, h, paragraphs: list[str], size=9,
                   color=None, align=PP_ALIGN.JUSTIFY):
    """Justified multi-paragraph text box."""
    if color is None:
        color = _DARK
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after  = Pt(5) if i < len(paragraphs) - 1 else Pt(0)
        run = p.add_run()
        run.text = para
        run.font.name  = "Inter"
        run.font.size  = Pt(size)
        run.font.color.rgb = color


def _embed(slide, buf: BytesIO, x, y, w, h):
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))


# ── Chart renderers ────────────────────────────────────────────────────────────

def _chart_all(df: pd.DataFrame,
               barrier_pct: float | None = None,
               cupon_annual: float | None = None,
               accrual_color: str = "#DC2626") -> BytesIO:
    """
    All underlyings, normalized 100 = inception.
    If barrier_pct + cupon_annual are provided, overlays the Range Accrual
    cumulative coupon line: each calendar day the worst-of sits above the
    barrier, it earns cupon_annual/252 of coupon (business-day approximation).
    """
    fig, ax = plt.subplots(figsize=(4.5, 3.5))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    for i, col in enumerate(df.columns):
        ax.plot(df.index, df[col], color=_LINE_COLORS[i % 4],
                linewidth=1.2, label=col)

    # ── Range Accrual coupon accumulation overlay ──────────────────────────────
    accrual = None
    if barrier_pct is not None and cupon_annual and not df.empty:
        worst = df.min(axis=1)
        daily_rate = cupon_annual / 252          # trading-day coupon rate
        above = (worst >= barrier_pct).astype(float)
        accrual = 100 + above.cumsum() * daily_rate
        ax.plot(accrual.index, accrual.values,
                color=accrual_color, linewidth=1.8, zorder=5,
                label=f"Cupón acumulado RA")

    ax.axhline(100, color="#9CA3AF", linewidth=0.8, linestyle="--", alpha=0.6)

    # Y-axis limits: include accrual line top so it's never clipped
    all_vals = df.values[~np.isnan(df.values)]
    if len(all_vals):
        data_min = float(all_vals.min())
        data_max = float(all_vals.max())
        if accrual is not None:
            data_max = max(data_max, float(accrual.max()))
        ymin = max(0.0, float((data_min - 10) // 5) * 5)
        ymax = float(np.ceil((data_max + 5) / 5.0) * 5)
        ax.set_ylim(ymin, ymax)

    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{v:.0f}%"))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%b-%y"))
    ax.xaxis.set_major_locator(mdates.MonthLocator(interval=2))
    plt.setp(ax.xaxis.get_majorticklabels(), rotation=45, ha="right", fontsize=9)
    plt.setp(ax.yaxis.get_majorticklabels(), fontsize=9)
    ax.spines[["top","right"]].set_visible(False)
    ax.spines[["left","bottom"]].set_color("#E5E7EB")
    ax.tick_params(colors="#6B7280", length=3)
    ax.grid(axis="y", linestyle="--", alpha=0.3, linewidth=0.5)
    ax.legend(fontsize=9, loc="upper left", framealpha=0.8,
              ncol=2 if len(df.columns) > 2 else 1,
              handlelength=1.5, borderpad=0.4, labelspacing=0.3)

    fig.tight_layout(pad=0.4)
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return buf


def _chart_worst(series: pd.Series) -> BytesIO:
    """Worst-of evolution only."""
    fig, ax = plt.subplots(figsize=(4.5, 3.2))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    ax.plot(series.index, series.values, color="#1A1A2E", linewidth=1.2)
    ax.axhline(100, color="#9CA3AF", linewidth=0.8, linestyle="--", alpha=0.6)

    vals = series.dropna().values
    if len(vals):
        ymin = max(0.0, float((vals.min() - 10) // 5) * 5)
        ymax = float(np.ceil((vals.max() + 5) / 5.0) * 5)
        ax.set_ylim(ymin, ymax)

    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{v:.0f}%"))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%b-%y"))
    ax.xaxis.set_major_locator(mdates.MonthLocator(interval=2))
    plt.setp(ax.xaxis.get_majorticklabels(), rotation=45, ha="right", fontsize=9)
    plt.setp(ax.yaxis.get_majorticklabels(), fontsize=9)
    ax.spines[["top","right"]].set_visible(False)
    ax.spines[["left","bottom"]].set_color("#E5E7EB")
    ax.tick_params(colors="#6B7280", length=3)
    ax.grid(axis="y", linestyle="--", alpha=0.3, linewidth=0.5)

    fig.tight_layout(pad=0.4)
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return buf


# ── Table builders ─────────────────────────────────────────────────────────────

def _cell_fmt(cell, text, size=9, bold=False, color=None,
              align=PP_ALIGN.LEFT, fill=None, bg=None):
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    if color is None:
        color = _DARK
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.05)
    tf.margin_right  = Inches(0.04)
    tf.margin_top    = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = str(text).split("\n")
    for li, line in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = p.space_after = Pt(0)
        run = p.add_run()
        run.text = str(line).strip()
        run.font.name  = "Inter"
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color


def _build_caract_table(slide, x, y, w, rows, pri: RGBColor):
    """Two-column characteristics table."""
    n = len(rows)
    if not n:
        return 0.0

    # Row height: 0.22" per text line
    line_counts = [max(1, str(v).count("\n") + 1) for _, v in rows]
    row_h_in = [0.22 * lc for lc in line_counts]
    total_h  = sum(row_h_in)

    tbl = slide.shapes.add_table(
        n, 2, Inches(x), Inches(y), Inches(w), Inches(total_h)
    ).table

    total_emu = int(Inches(w))
    tbl.columns[0].width = int(total_emu * 0.44)
    tbl.columns[1].width = int(total_emu * 0.56)

    for ri, rh in enumerate(row_h_in):
        tbl.rows[ri].height = int(Inches(rh))

    for ri, (label, value) in enumerate(rows):
        bg = _LIGHT if ri % 2 == 0 else _WHITE
        _cell_fmt(tbl.cell(ri, 0), label, bold=True,  align=PP_ALIGN.LEFT,  bg=bg)
        _cell_fmt(tbl.cell(ri, 1), value, bold=False, align=PP_ALIGN.RIGHT, bg=bg)

    return total_h


def _build_summary_table(slide, x, y, w, h, name, fi, fv, fa,
                         rendimientos, worst_u, worst_v,
                         rend_producto, last_pmt, pri: RGBColor):
    """Full-width horizontal summary table."""
    rend_lines = "\n".join(f"{u}: {v:+.2f}%" for u, v in rendimientos.items())
    worst_str  = f"{worst_u}: {worst_v:+.2f}%" if worst_u else "—"
    rend_p_str = f"{rend_producto:.3f}%" if rend_producto is not None else "—"
    last_pmt_str = _fmt_short(last_pmt)

    headers = [
        "Producto",
        "Fecha\nInicio",
        "Fecha\nVencimiento",
        "Fecha de\nAutocall",
        f"Pago cupón\n({last_pmt_str})",
        "Rendimiento\nsubyacente",
        "Rend. del peor\nsubyacente",
        "Rendimiento\ndel Producto",
    ]
    data = [
        name,
        _fmt_short(fi),
        _fmt_short(fv),
        _fmt_short(fa),
        "—",
        rend_lines or "—",
        worst_str,
        rend_p_str,
    ]
    col_pcts = [0.195, 0.085, 0.105, 0.105, 0.090, 0.175, 0.130, 0.115]
    n = len(headers)

    tbl = slide.shapes.add_table(2, n, Inches(x), Inches(y),
                                  Inches(w), Inches(h)).table

    total_emu = int(Inches(w))
    for ci, pct in enumerate(col_pcts):
        tbl.columns[ci].width = int(total_emu * pct)

    tbl.rows[0].height = int(Inches(h * 0.38))
    tbl.rows[1].height = int(Inches(h * 0.62))

    for ci, (hdr, dat) in enumerate(zip(headers, data)):
        _cell_fmt(tbl.cell(0, ci), hdr, size=7, bold=True,
                  color=_WHITE, align=PP_ALIGN.CENTER, fill=pri)
        _cell_fmt(tbl.cell(1, ci), dat, size=8, bold=(ci == 0),
                  align=PP_ALIGN.CENTER, bg=_WHITE)


def _build_coupon_table(slide, x, y, w, coupon_rows, total_str, pri: RGBColor):
    """N° | Fecha de Pago | Cupón table."""
    n_data = len(coupon_rows)
    n_rows = 1 + n_data + 1   # header + data + total
    row_h  = 0.27
    total_h = row_h * n_rows

    tbl = slide.shapes.add_table(
        n_rows, 3, Inches(x), Inches(y), Inches(w), Inches(total_h)
    ).table

    emu_w = int(Inches(w))
    tbl.columns[0].width = int(emu_w * 0.18)
    tbl.columns[1].width = int(emu_w * 0.45)
    tbl.columns[2].width = int(emu_w * 0.37)

    rh_emu = int(Inches(row_h))
    for ri in range(n_rows):
        tbl.rows[ri].height = rh_emu

    # Header
    for ci, hdr in enumerate(["N° de Pago", "Fecha de Pago", "Cupón"]):
        _cell_fmt(tbl.cell(0, ci), hdr, bold=True,
                  color=_WHITE, align=PP_ALIGN.CENTER, fill=pri)

    # Data rows
    for ri, (n, fecha, cupon_s) in enumerate(coupon_rows, start=1):
        bg = _LIGHT if ri % 2 == 0 else _WHITE
        for ci, val in enumerate([str(n), fecha, cupon_s]):
            _cell_fmt(tbl.cell(ri, ci), val, align=PP_ALIGN.CENTER, bg=bg)

    # Total row
    total_ri = n_data + 1
    for ci, val in enumerate(["Total pagado", "", total_str]):
        _cell_fmt(tbl.cell(total_ri, ci), val, bold=True,
                  color=_WHITE, align=PP_ALIGN.CENTER, fill=pri)

    return total_h


# ── Narrative builder ──────────────────────────────────────────────────────────

def _narrative(name, event_type, ac_date, fi, fv, ac_idx,
               cupon, barrera_pct, ganancia_max, factor_part,
               rend_producto) -> list[str]:
    ki = 100 - barrera_pct if barrera_pct else None
    ki_str = f"{ki:.0f}%" if ki else "—"
    cupon_str = _pct(cupon)

    if event_type == "Autocall":
        ord_es = _ORDINAL_ES[ac_idx - 1] if ac_idx and 1 <= ac_idx <= 10 else f"#{ac_idx}"
        rend_s = f"{rend_producto:.3f}%" if rend_producto else "—"
        p1 = (
            f"El {name} fue llamado anticipadamente en la {ord_es} fecha de observación "
            f"de autocall y alcanzó una rentabilidad de {rend_s} para el período comprendido "
            f"entre el {_fmt_long(fi)} y el {_fmt_long(ac_date)}. El producto brindaba la "
            f"posibilidad de obtener un cupón anual contingente de {cupon_str} si en las "
            f"observaciones el subyacente estaba por encima del {ki_str} de su nivel inicial."
        )
        paras = [p1]
        if factor_part:
            fp = factor_part * 100 if factor_part <= 2 else factor_part
            paras.append(
                f"Si a vencimiento el subyacente se encontraba por debajo del {ki_str} de su "
                f"nivel inicial, el inversionista habría estado expuesto en un {fp:.0f}% "
                f"al rendimiento negativo del subyacente."
            )
        return paras

    if event_type == "Vencimiento":
        gmax = ganancia_max
        gmax_s = f"{gmax*100:.2f}%" if gmax and gmax <= 1 else (f"{gmax:.2f}%" if gmax else "—")
        return [
            f"El {name} llegó a vencimiento obteniendo una rentabilidad de {gmax_s} "
            f"para el período comprendido entre el {_fmt_long(fi)} y el {_fmt_long(fv)}. "
            f"El producto brindaba la posibilidad de obtener un cupón anual contingente "
            f"de {cupon_str} si en las observaciones el subyacente estaba por encima "
            f"del {ki_str} de su nivel inicial."
        ]

    return [
        f"El {name} es un producto estructurado que ofrece un cupón anual contingente "
        f"de {cupon_str} si el subyacente se mantiene por encima del {ki_str} de su "
        f"nivel inicial en las fechas de observación."
    ]


# ── Price data ─────────────────────────────────────────────────────────────────

def _fetch_perf(unds, fecha_inicio, cutoff):
    if not unds or not fecha_inicio or not YF_OK:
        return pd.DataFrame(), pd.Series(dtype=float)
    yf_tickers = [resolve_ticker(u) for u in unds if resolve_ticker(u)]
    if not yf_tickers:
        return pd.DataFrame(), pd.Series(dtype=float)
    try:
        raw = yf.download(
            yf_tickers,
            start=str(fecha_inicio - timedelta(days=7)),
            end=str((cutoff or date.today()) + timedelta(days=2)),
            auto_adjust=False, progress=False,
        )
        if raw.empty:
            return pd.DataFrame(), pd.Series(dtype=float)
        if isinstance(raw.columns, pd.MultiIndex):
            closes = raw["Close"].copy()
        else:
            closes = raw[["Close"]].copy()
            closes.columns = yf_tickers[:1]
        closes.index = pd.to_datetime(closes.index).tz_localize(None)
        pre = closes[closes.index.date <= fecha_inicio]
        if pre.empty:
            pre = closes[closes.index.date <= fecha_inicio + timedelta(days=5)]
        if pre.empty:
            return pd.DataFrame(), pd.Series(dtype=float)
        init = pre.iloc[-1]
        norm = closes / init * 100
        norm.columns = unds[:len(norm.columns)]
        return norm, norm.min(axis=1)
    except Exception:
        return pd.DataFrame(), pd.Series(dtype=float)


# ══════════════════════════════════════════════════════════════════════════════
# PARTICIPATION PRODUCT FACTSHEETS
# Dual Directional (Twin Win), Call Spread, Airbag, Capital Protected
# ══════════════════════════════════════════════════════════════════════════════

_PART_ELEM_TYPES = {
    "Long Call (ATM)", "Long Call (ITM)",
    "Capital Protected Participation",
}

def is_participation_product(product: dict) -> bool:
    return str(product.get("elemento_1_tipo") or "").strip() in _PART_ELEM_TYPES


def _chart_payoff_dd(barrier_pct: float,
                     factor_part: float = 1.0,
                     ganancia_max_str: str = "Ilimitada") -> BytesIO:
    """
    Twin Win / Dual Directional payoff diagram.
    barrier_pct : BUFFER size as a percentage  (e.g. 16.25 for a 16.25 % buffer)
                  The actual barrier LEVEL = 100 % – barrier_pct  (e.g. 83.75 %)
    Uses underlying LEVEL on x-axis, product LEVEL on y-axis (both % of initial).
    """
    blevel = 100 - barrier_pct          # barrier level   e.g. 83.75
    x_min, x_max = 55.0, 148.0

    # --- three segments ----------------------------------------------------
    x1 = np.linspace(x_min, blevel, 300)     # loss zone  (x < blevel)
    y1 = x1.copy()                             # product = underlying

    x2 = np.linspace(blevel, 100.0, 300)      # twin-win zone
    y2 = 200.0 - x2                            # product = 200 – underlying

    x3 = np.linspace(100.0, x_max, 300)       # upside zone
    y3 = 100.0 + (x3 - 100.0) * factor_part

    fig, ax = plt.subplots(figsize=(4.2, 3.6))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    ax.plot(x1, y1, color="#1A1A2E", linewidth=2.2)
    # jump — show as short dashed vertical at the barrier
    ax.plot([blevel, blevel], [blevel, 100 + barrier_pct],
            color="#1A1A2E", linewidth=1.4, linestyle=":")
    ax.plot(x2, y2, color="#1A1A2E", linewidth=2.2)
    ax.plot(x3, y3, color="#1A1A2E", linewidth=2.2)

    # subtle shading for twin-win buffer zone
    ax.fill_between(x2, 100, y2, alpha=0.10, color="#2563EB")

    # reference lines
    ax.axhline(100, color="#9CA3AF", linewidth=0.7, linestyle="--", alpha=0.5)
    ax.axvline(100, color="#9CA3AF", linewidth=0.7, linestyle="--", alpha=0.5)
    ax.axvline(blevel, color="#DC2626", linewidth=0.7, linestyle="--", alpha=0.45)

    # barrier level label on x-axis
    ax.annotate(f"{blevel:.2f}%",
                xy=(blevel, x_min + 0.5), fontsize=7.5,
                ha="center", va="bottom", color="#DC2626", fontweight="bold")

    # "Punto inicial" dot + annotation
    ax.scatter([100], [100], color="#1A1A2E", s=18, zorder=6)
    ax.annotate("Punto inicial\n(100%)", xy=(100, 100),
                xytext=(104, 91), fontsize=7.5, ha="left", va="top",
                color="#475569",
                arrowprops=dict(arrowstyle="->", color="#9CA3AF", lw=0.7))

    # max-gain label with arrow
    gmax_low = ganancia_max_str.lower()
    if "ilimitad" in gmax_low or "unlimited" in gmax_low:
        label = "Ganancia máxima:\nIlimitada"
    else:
        label = f"Ganancia máxima:\n{ganancia_max_str}"
    arrow_x = x_max - 5
    arrow_y = y3[-1]
    ax.annotate(label,
                xy=(arrow_x, arrow_y - 3), fontsize=7.5,
                ha="center", va="top", color="#1A1A2E",
                arrowprops=None)
    ax.annotate("", xy=(arrow_x + 3, arrow_y + 3),
                xytext=(arrow_x, arrow_y - 3),
                arrowprops=dict(arrowstyle="->", color="#1A1A2E", lw=1.0))

    # slope labels
    ax.text(68, 65, "1:1", fontsize=7.5, ha="center", color="#6B7280")
    ax.text(90, 112, "1:1", fontsize=7.5, ha="center", color="#6B7280")
    fp_label = f"{factor_part:.2f}x" if factor_part != 1.0 else "1:1"
    ax.text(132, 135, fp_label, fontsize=7.5, ha="center", color="#6B7280")

    ax.set_xlabel("Subyacente", fontsize=8.5, color="#6B7280")
    ax.set_ylabel("Producto",   fontsize=8.5, color="#6B7280")
    ax.set_xlim(x_min, x_max)
    ax.set_ylim(x_min, x_max)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{v:.0f}%"))
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{v:.0f}%"))
    ax.spines[["top","right"]].set_visible(False)
    ax.spines[["left","bottom"]].set_color("#E5E7EB")
    ax.tick_params(colors="#6B7280", length=3, labelsize=7.5)
    ax.grid(axis="both", linestyle="--", alpha=0.15, linewidth=0.5)

    fig.tight_layout(pad=0.35)
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return buf


def _scen_returns(barrier_pct: float, n_unds: int) -> list[tuple]:
    """
    Generate 5 illustrative scenarios for the scenario table.
    Returns list of (und_rets: list[float], worst: float, prod_ret: float).
    und_rets are percentage returns (e.g., 45.0 for 45%).
    """
    b = barrier_pct  # e.g. 16.25
    # 5 scenario worst-of targets (as %)
    targets = [45.0, 18.0, -(b * 0.31), -(b - 0.5), -(b + 19.0)]

    # Build illustrative per-underlying returns for up to 3 underlyings
    # such that min(row) == target
    raw3 = [
        [45.0, 55.0, 60.0],
        [30.0, 18.0, 20.0],
        [10.0, 7.0,  targets[2]],
        [-b*0.31, -b*0.62, targets[3]],
        [-(b+9.0), -(b+5.0), -(b+19.0)],
    ]

    out = []
    for i, row in enumerate(raw3):
        und_rets = row[:n_unds]
        worst    = min(und_rets)
        w        = worst / 100
        if w >= 0:
            p_ret = w * 1.0  # factor_part applied later via caller
        elif w >= -b/100:
            p_ret = -w
        else:
            p_ret = w
        out.append((und_rets, worst, p_ret * 100))
    return out


def _build_scenario_tbl(slide, x, y, w, unds: list, barrier_pct: float,
                         factor_part: float, pri: RGBColor):
    """Scenario example table (rows = underlyings + worst-of + product return)."""
    n_u = min(len(unds), 3)
    scenarios = _scen_returns(barrier_pct, n_u)

    row_labels = unds[:n_u] + ["Rend. subyacente", "Rentabilidad\ndel Producto"]
    n_rows = len(row_labels) + 1  # +1 for header row
    n_cols = 1 + 5  # label + 5 examples
    row_h  = 0.21
    total_h = n_rows * row_h

    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(total_h)
    ).table

    emu_w = int(Inches(w))
    tbl.columns[0].width = int(emu_w * 0.30)
    for ci in range(1, 6):
        tbl.columns[ci].width = int(emu_w * 0.14)

    for ri in range(n_rows):
        tbl.rows[ri].height = int(Inches(row_h))

    # Header row
    _cell_fmt(tbl.cell(0, 0), "Subyacente",  size=7.5, bold=True,
              align=PP_ALIGN.CENTER, fill=pri, color=_WHITE)
    for ci in range(1, 6):
        _cell_fmt(tbl.cell(0, ci), f"Ejm {ci}", size=7.5, bold=True,
                  align=PP_ALIGN.CENTER, fill=pri, color=_WHITE)

    # Underlying rows
    for ri in range(1, n_u + 1):
        und = unds[ri - 1]
        bg = _LIGHT if ri % 2 == 0 else _WHITE
        _cell_fmt(tbl.cell(ri, 0), und, size=7.5, bold=True, bg=bg)
        for ci, (und_rets, _, _) in enumerate(scenarios, 1):
            pct_val = und_rets[ri - 1] if ri - 1 < len(und_rets) else 0.0
            sign = "+" if pct_val >= 0 else ""
            _cell_fmt(tbl.cell(ri, ci), f"{sign}{pct_val:.0f}%",
                      size=7.5, align=PP_ALIGN.CENTER, bg=bg)

    # Worst-of row
    ri_worst = n_u + 1
    bg = _LIGHT if ri_worst % 2 == 0 else _WHITE
    _cell_fmt(tbl.cell(ri_worst, 0), "Rend. subyacente",
              size=7.5, bold=True, bg=bg)
    for ci, (_, worst, _) in enumerate(scenarios, 1):
        sign = "+" if worst >= 0 else ""
        _cell_fmt(tbl.cell(ri_worst, ci), f"{sign}{worst:.2f}%",
                  size=7.5, align=PP_ALIGN.CENTER, bg=bg)

    # Product return row (bold, color-coded)
    ri_prod = n_u + 2
    bg = _LIGHT if ri_prod % 2 == 0 else _WHITE
    _cell_fmt(tbl.cell(ri_prod, 0), "Rentabilidad\ndel Producto",
              size=7.5, bold=True, bg=bg)
    for ci, (_, _, p_ret) in enumerate(scenarios, 1):
        # Apply factor_part to positive returns only
        adj = p_ret * factor_part if p_ret > 0 else p_ret
        sign = "+" if adj >= 0 else ""
        col = _rgb("#16A34A") if adj > 0 else (_DARK if adj == 0 else _rgb("#DC2626"))
        _cell_fmt(tbl.cell(ri_prod, ci), f"{sign}{adj:.2f}%",
                  size=7.5, bold=True, align=PP_ALIGN.CENTER, bg=bg, color=col)

    return total_h


def _is_dual_directional(product: dict) -> bool:
    return str(product.get("elemento_2_tipo") or "").strip() == "KO Put (ATM)"


def _narrative_participation(product: dict, elem_tipo: str) -> list[str]:
    """Narrative description paragraph for participation products."""
    asset_class  = str(product.get("asset_class") or "renta variable")
    plazo_meses  = product.get("plazo_meses")
    plazo_str    = f"{int(float(plazo_meses))} meses" if plazo_meses else "—"
    factor_part  = _sf(product.get("factor_participacion")) or 1.0
    barrera_cap  = _sf(product.get("barrera_capital")) or 0
    buffer_pct   = (1 - barrera_cap) * 100 if barrera_cap <= 1 else (100 - barrera_cap)

    ganancia_max = str(product.get("ganancia_maxima") or "Ilimitada").strip()
    if ganancia_max.lower() in ("ilimitada", "unlimited", ""):
        gmax_phrase = "sin límite de ganancia"
    else:
        gmax_phrase = f"hasta un máximo de {ganancia_max}"

    fp_str = f"{factor_part:.2f}x" if factor_part else "1.00x"

    if _is_dual_directional(product):
        return [
            f"El producto brinda exposición a la {asset_class.lower()} a {plazo_str}. "
            f"El inversionista se beneficiará del upside del subyacente multiplicado por "
            f"un factor de participación de {fp_str} {gmax_phrase}. A su vez, se "
            f"beneficiará 1:1 del downside hasta una caída de {buffer_pct:.2f}%. "
            f"Luego de ese límite, el inversionista está expuesto en un 100% al "
            f"rendimiento negativo del subyacente."
        ]
    if elem_tipo in ("Long Call (ATM)", "Long Call (ITM)"):
        return [
            f"El producto brinda exposición a la {asset_class.lower()} a {plazo_str}. "
            f"El inversionista participa del upside del subyacente con un factor de "
            f"{fp_str} {gmax_phrase}. El capital está protegido hasta una caída de "
            f"{buffer_pct:.2f}% del nivel inicial."
        ]
    # Capital Protected default
    return [
        f"El producto brinda exposición a la {asset_class.lower()} a {plazo_str}. "
        f"El inversionista participa del upside del subyacente con un factor de "
        f"{fp_str} {gmax_phrase}. El 100% del capital está protegido al vencimiento "
        f"independientemente del comportamiento del subyacente."
    ]


def _scenario_bullets(product: dict, buffer_pct: float) -> list[str]:
    """Bullet scenario descriptions for the right column."""
    b_str = f"{buffer_pct:.2f}%"
    if _is_dual_directional(product):
        return [
            "A vencimiento:",
            f"▪  Escenario 1: Si el subyacente tiene un rendimiento positivo, el "
            f"inversionista recibe la apreciación del subyacente multiplicado por el "
            f"factor de participación.",
            f"▪  Escenario 2: Si el subyacente tiene un rendimiento negativo y mayor a "
            f"-{b_str}, el inversionista se beneficia 1:1 de la pérdida del subyacente.",
            f"▪  Escenario 3: Si el subyacente tiene un rendimiento negativo y menor a "
            f"-{b_str}, el inversionista experimenta el mismo rendimiento del subyacente.",
        ]
    return [
        "A vencimiento:",
        f"▪  Escenario 1: Rendimiento positivo — el inversionista recibe la apreciación "
        f"del subyacente multiplicado por el factor de participación.",
        f"▪  Escenario 2: Caída dentro del buffer ({b_str}) — el capital queda protegido.",
        f"▪  Escenario 3: Caída por debajo de {b_str} — exposición al rendimiento "
        f"negativo del subyacente.",
    ]


def generate_factsheet_participation(
    product: dict,
    event_type: str,
    company_name: str,
    primary: str,
    secondary: str = "#DC2626",
    logo_bytes: bytes | None = None,
    disclaimer: str | None = None,
) -> bytes:
    """
    A4 factsheet for participation products (Dual Directional, Call Spread, Airbag…).
    Layout mirrors the Credicorp Capital twin-win template:
      top-left : narrative + initial prices table
      top-right: characteristics table
      bot-left : payoff diagram
      bot-right: scenario bullets + scenario table
    """
    if not _PPTX:
        raise ImportError("python-pptx not installed.")

    PRI = _rgb(primary)
    SEC = _rgb(secondary)

    # ── Extract product fields ─────────────────────────────────────────────
    name        = str(product.get("nombre_producto") or "")
    asset_class = str(product.get("asset_class") or "—")
    tipo        = str(product.get("tipo") or "Nota Estructurada")
    moneda      = str(product.get("moneda") or "USD")
    isin        = str(product.get("isin") or "—")
    contraparte = str(product.get("contraparte") or "—")
    vehiculo    = str(product.get("vehiculo") or "—")
    tipo_cliente = str(product.get("tipo_cliente") or "—")

    fecha_inicio  = _parse_date(product.get("fecha_inicio") or product.get("fecha_strike"))
    fecha_obs_fin = _parse_date(product.get("fecha_obs_final"))
    fecha_vencto  = _parse_date(
        product.get("fecha_vencimiento") or
        product.get("fecha_obs_final") or
        product.get("fecha_obs_final_ac")
    )

    plazo_meses  = _sf(product.get("plazo_meses"))
    factor_part  = _sf(product.get("factor_participacion")) or 1.0
    barrera_cap  = _sf(product.get("barrera_capital")) or 0
    if barrera_cap > 1:
        barrera_cap /= 100
    buffer_pct   = (1.0 - barrera_cap) * 100   # e.g. 16.25

    ganancia_max = str(product.get("ganancia_maxima") or "Ilimitada").strip()
    elem_tipo    = str(product.get("elemento_1_tipo") or "")

    unds = [str(product.get(f"underlying_{i}")).strip()
            for i in range(1, 5)
            if product.get(f"underlying_{i}") and
               str(product.get(f"underlying_{i}")).strip() not in ("", "nan", "None")]

    strikes = {u: _sf(product.get(f"strike_{i}"))
               for i, u in enumerate(unds, 1)
               if _sf(product.get(f"strike_{i}"))}

    und_long = [_UND_LABELS.get(u, u) for u in unds]
    if len(unds) > 1:
        und_desc = "Worst of: " + ", ".join(und_long)
        sub_desc = ("Opción Worst of sobre el " + " / ".join(und_long))
    else:
        und_desc = und_long[0] if und_long else "—"
        sub_desc = und_desc

    # ── Derived strings ────────────────────────────────────────────────────
    plazo_str   = f"{int(plazo_meses)} meses" if plazo_meses else "—"
    moneda_str  = "Dólares Americanos" if "USD" in moneda.upper() else moneda
    fp_str      = f"{factor_part:.2f}x"
    gmax_disp   = ganancia_max if ganancia_max else "Ilimitada"
    barr_disp   = f"{buffer_pct:.2f}%"  # show buffer size (e.g., "16.25%")

    # ── Characteristics rows ───────────────────────────────────────────────
    caract_rows = [
        ("Plataforma",           vehiculo),
        ("Cliente",              tipo_cliente),
        ("Clasificación",        asset_class),
        ("Tipo de producto",     tipo),
        ("Concentración máxima", "5% de la cartera"),
        ("Subyacente",           sub_desc),
        ("Moneda",               moneda_str),
        ("Plazo",                plazo_str),
        ("Riesgo del producto",  contraparte),
        ("Pago del Retorno",     "Al vencimiento"),
        ("Barrera",              barr_disp),
        ("Factor de Participación", fp_str),
        ("Ganancia Máxima",      gmax_disp),
        ("ISIN",                 isin),
        ("Fecha de Observación inicial", _fmt_short(fecha_inicio)),
        ("Fecha de Observación final",   _fmt_short(fecha_obs_fin or fecha_vencto)),
        ("Fecha de Vencimiento (pago)",  _fmt_short(fecha_vencto)),
    ]
    caract_rows = [(k, v) for k, v in caract_rows if v and v not in ("—", "")]

    # ── Narrative & bullets ────────────────────────────────────────────────
    narrative = _narrative_participation(product, elem_tipo)
    bullets   = _scenario_bullets(product, buffer_pct)

    # ── Charts ────────────────────────────────────────────────────────────
    payoff_buf = _chart_payoff_dd(buffer_pct, factor_part, ganancia_max)

    # ── BUILD SLIDE ────────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(_W)
    prs.slide_height = Inches(_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    # Y constants
    Y_HDR   = 0.00;  H_HDR  = 0.55
    Y_TTL   = 0.58;  H_TTL  = 0.26
    Y_INFO  = 0.88;  H_INFO = 0.22
    Y_BARS  = 1.14;  H_BAR  = 0.22
    Y_TOP   = 1.40;  H_TOP  = 4.00
    Y_BBAR  = 5.44;  H_BBAR = 0.22
    Y_BOT   = 5.70;  H_BOT  = 4.60
    Y_DISC  = 10.34; H_DISC = 0.60
    Y_FOOT  = 10.96; H_FOOT = 0.22

    # 1. Header bar + logo/name
    _rect(slide, 0, Y_HDR, _W, H_HDR, PRI)
    if logo_bytes:
        slide.shapes.add_picture(
            BytesIO(logo_bytes), Inches(_ML), Inches(Y_HDR + 0.04),
            width=None, height=Inches(H_HDR - 0.08))
    else:
        _txt(slide, _ML, Y_HDR + 0.09, 2.5, H_HDR - 0.18,
             company_name, size=9, bold=True, color=_WHITE)

    bw = 1.0
    badge = _rect(slide, _W - bw - 0.20, Y_HDR + 0.09, bw, H_HDR - 0.18, _WHITE)
    btf = badge.text_frame
    btf.margin_left = btf.margin_right = Inches(0.05)
    btf.margin_top  = btf.margin_bottom = Inches(0.02)
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = event_type.upper()
    br.font.name = "Inter"; br.font.size = Pt(9)
    br.font.bold = True; br.font.color.rgb = PRI

    # 2. Product title
    _txt(slide, _ML, Y_TTL, _CW, H_TTL,
         f"{event_type.upper()} – {name}",
         size=11, bold=True, color=SEC)

    # 3. Full-width INFO bar
    _bar(slide, _ML, Y_INFO, _CW, H_INFO, "INFORMACIÓN DEL PRODUCTO", SEC, size=9)

    # 4. Dual section bars
    _bar(slide, _ML, Y_BARS, _LW, H_BAR, "DESCRIPCIÓN Y ESTRATEGIA", PRI, size=9)
    _bar(slide, _RX, Y_BARS, _RW, H_BAR, "CARACTERÍSTICAS GENERALES",  PRI, size=9)

    # 5. LEFT TOP — narrative + initial prices
    cursor = Y_TOP
    _multiline_txt(slide, _ML, cursor, _LW, 1.30, narrative, size=9,
                   align=PP_ALIGN.JUSTIFY)
    cursor += 1.38

    # initial prices table
    if strikes:
        n_pr  = len(strikes)
        pr_h  = 0.22 * (n_pr + 1)  # header + data rows
        tbl_p = slide.shapes.add_table(
            n_pr + 1, 2,
            Inches(_ML), Inches(cursor),
            Inches(_LW), Inches(pr_h)
        ).table
        emu = int(Inches(_LW))
        tbl_p.columns[0].width = int(emu * 0.42)
        tbl_p.columns[1].width = int(emu * 0.58)
        for ri in range(n_pr + 1):
            tbl_p.rows[ri].height = int(Inches(0.22))
        _cell_fmt(tbl_p.cell(0, 0), "Subyacente", bold=True,
                  fill=PRI, color=_WHITE, align=PP_ALIGN.CENTER, size=8)
        _cell_fmt(tbl_p.cell(0, 1), "Precio Inicial", bold=True,
                  fill=PRI, color=_WHITE, align=PP_ALIGN.CENTER, size=8)
        for ri, (u, val) in enumerate(strikes.items(), 1):
            bg = _LIGHT if ri % 2 == 0 else _WHITE
            _cell_fmt(tbl_p.cell(ri, 0), u, size=8, bold=True,
                      align=PP_ALIGN.CENTER, bg=bg)
            _cell_fmt(tbl_p.cell(ri, 1), f"{val:,.2f}", size=8,
                      align=PP_ALIGN.CENTER, bg=bg)
        cursor += pr_h + 0.12

    # small footnote below prices
    _txt(slide, _ML, cursor, _LW, 0.28,
         "*Rating Crediticio: Grado de Inversión. Las notas estructuradas califican "
         "como \"senior unsecured debt\" en la estructura de capital de las contrapartes.",
         size=7.5, color=_MUTED, italic=True)

    # 6. RIGHT TOP — characteristics table
    _build_caract_table(slide, _RX, Y_TOP, _RW, caract_rows, PRI)

    # 7. Bottom dual bars
    _bar(slide, _ML, Y_BBAR, _LW, H_BBAR, "ESTRUCTURA DEL PRODUCTO", SEC, size=9)
    _bar(slide, _RX, Y_BBAR, _RW, H_BBAR, "ESCENARIOS DE CAPITAL",  SEC, size=9)

    # 8. LEFT BOTTOM — payoff diagram
    _embed(slide, payoff_buf, _ML, Y_BOT, _LW, H_BOT * 0.78)

    # 9. RIGHT BOTTOM — scenario bullets + table
    blt_x = _RX; blt_y = Y_BOT; blt_w = _RW
    # bullets text
    tb = slide.shapes.add_textbox(
        Inches(blt_x), Inches(blt_y), Inches(blt_w), Inches(1.50))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top  = tf.margin_bottom = Inches(0)
    for li, line in enumerate(bullets):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.space_before = Pt(3) if li > 0 else Pt(0)
        p.space_after  = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name  = "Inter"
        run.font.size  = Pt(8.5)
        run.font.bold  = (li == 0)
        run.font.color.rgb = _DARK

    # scenario sub-bar
    scen_bar_y = Y_BOT + 1.60
    _bar(slide, _RX, scen_bar_y, _RW, 0.20,
         "TABLA DE EJEMPLO DE RENDIMIENTOS", PRI, size=8)

    # scenario table
    scen_tbl_y = scen_bar_y + 0.24
    _build_scenario_tbl(
        slide, _RX, scen_tbl_y, _RW,
        unds, buffer_pct, factor_part, PRI
    )

    # 10. Disclaimer / footnote
    disc_text = (
        disclaimer if disclaimer
        else "La colocación de la nota estructurada es por oferta privada. "
             "Credicorp Capital no garantiza el pago de rendimiento y/o principal "
             "de los valores. Los factores de riesgo incluyen riesgo de mercado, "
             "riesgo del emisor y riesgo de liquidez."
    )
    _txt(slide, _ML, Y_DISC, _CW, H_DISC, disc_text,
         size=7.5, italic=True, color=_MUTED)

    _txt(slide, _ML, Y_FOOT, _CW, H_FOOT,
         f"*Al vencimiento. La rentabilidad esperada no incluye el efecto del impuesto "
         f"correspondiente a cada inversionista.",
         size=7.5, italic=True, color=_MUTED)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ── Main entry point ───────────────────────────────────────────────────────────

def generate_factsheet_pdf(product: dict, event_type: str, company_name: str,
                           primary: str, secondary: str = "#DC2626",
                           verified_autocall_date=None,
                           logo_bytes: bytes | None = None,
                           disclaimer: str | None = None) -> bytes:
    """
    Returns PPTX bytes (A4 portrait).
    primary    → header bar, sub-bars, table headers
    secondary  → main section bars (DETALLE, EVOLUCIÓN) and title
    logo_bytes → optional PNG/JPG bytes; embedded in the left side of the header bar
    disclaimer → optional text shown at the bottom instead of the default footnote
    """
    if not _PPTX:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    PRI = _rgb(primary)
    SEC = _rgb(secondary)

    # ── Extract product fields ─────────────────────────────────────────────────
    name        = str(product.get("nombre_producto") or "")
    asset_class = str(product.get("asset_class") or product.get("tipo") or "—")
    moneda      = str(product.get("moneda") or "USD")
    isin        = str(product.get("isin") or "—")
    estrategia  = str(product.get("estrategia") or product.get("tipo") or "—")

    fecha_inicio = _parse_date(product.get("fecha_inicio") or product.get("fecha_strike"))
    fecha_vencto = _parse_date(
        product.get("fecha_vencimiento") or
        product.get("fecha_obs_final") or
        product.get("fecha_obs_final_ac")
    )

    plazo_meses  = _sf(product.get("plazo_meses"))
    cupon_cont   = _sf(product.get("cupon_contingente"))
    cupon_fijo   = _sf(product.get("cupon_fijo"))
    cupon        = cupon_cont or cupon_fijo

    barrera_cap = _sf(product.get("barrera_capital")) or 0
    barrera_pct = barrera_cap * 100 if barrera_cap <= 1 else barrera_cap  # e.g. 25.0
    barrera_ki  = 100 - barrera_pct                                        # e.g. 75.0

    # Coupon barrier for Range Accrual overlay (may differ from capital barrier)
    _bc_raw = _sf(product.get("barrera_cupon"))
    if _bc_raw:
        barrier_cupon = _bc_raw * 100 if _bc_raw <= 1 else _bc_raw
    else:
        barrier_cupon = barrera_ki  # fall back to capital barrier level

    trigger     = _sf(product.get("trigger_autocall"))
    ganancia_max = _sf(product.get("ganancia_maxima"))
    factor_part  = _sf(product.get("factor_participacion"))

    periodo_nc   = str(product.get("periodo_sin_autocall") or "—")
    obs_autocall = str(product.get("observacion_autocall") or "Trimestral")
    obs_cupon    = str(product.get("observacion_cupon")    or "Trimestral")

    # Underlyings
    unds = [str(product.get(f"underlying_{i}")).strip()
            for i in range(1, 5)
            if product.get(f"underlying_{i}") and
               str(product.get(f"underlying_{i}")).strip() not in ("", "nan", "None")]

    strikes = {u: _sf(product.get(f"strike_{i}"))
               for i, u in enumerate(unds, 1)
               if _sf(product.get(f"strike_{i}"))}
    spots   = {u: _sf(product.get(f"spot_{i}"))
               for i, u in enumerate(unds, 1)
               if _sf(product.get(f"spot_{i}"))}

    # Autocall observation dates
    ac_dates = [d for i in range(1, 11)
                if (d := _parse_date(product.get(f"fecha_autocall_{i}")))]

    # Resolve actual autocall date
    actual_ac = None
    ac_idx    = None
    if event_type == "Autocall":
        actual_ac = _parse_date(verified_autocall_date)
        if not actual_ac and ac_dates:
            past = [d for d in ac_dates if d <= date.today()]
            actual_ac = past[-1] if past else ac_dates[0]
        if actual_ac and actual_ac in ac_dates:
            ac_idx = ac_dates.index(actual_ac) + 1

    # ── Price performance ──────────────────────────────────────────────────────
    cutoff = actual_ac or fecha_vencto or date.today()
    perf_df, worst_series = _fetch_perf(unds, fecha_inicio, cutoff)

    # ── Rendimientos ──────────────────────────────────────────────────────────
    rendimientos: dict[str, float] = {}
    worst_u, worst_v = None, None
    for u in unds:
        s, sp = strikes.get(u), spots.get(u)
        if s and sp and s > 0:
            r = (sp / s - 1) * 100
            rendimientos[u] = r
            if worst_v is None or r < worst_v:
                worst_u, worst_v = u, r

    # ── Coupon per period (uses actual observation frequency, not hardcoded /4) ──
    cupon_annual = (cupon * 100 if cupon and cupon <= 1 else cupon) if cupon else 0
    _FREQ_MAP = {"mensual": 12, "bimensual": 6, "trimestral": 4, "semestral": 2, "anual": 1}
    _freq = _FREQ_MAP.get(obs_autocall.lower().strip(), 4)
    cupon_period = cupon_annual / _freq if cupon_annual and _freq else 0

    # ── Coupon payment schedule ────────────────────────────────────────────────
    cutoff_coupons = actual_ac or date.today()
    coupon_rows = []
    for i, obs_d in enumerate(ac_dates, 1):
        if obs_d > cutoff_coupons:
            break
        pmt = obs_d + timedelta(days=15)
        while pmt.weekday() >= 5:
            pmt += timedelta(days=1)
        coupon_rows.append((i, _fmt_short(pmt), f"{cupon_period:.3f}%"))

    total_coupon_str = f"{len(coupon_rows) * cupon_period:.3f}%"

    last_pmt = None
    if coupon_rows and ac_dates:
        idx = len(coupon_rows) - 1
        if idx < len(ac_dates):
            p = ac_dates[idx] + timedelta(days=15)
            while p.weekday() >= 5:
                p += timedelta(days=1)
            last_pmt = p

    # ── Product return ─────────────────────────────────────────────────────────
    # = total coupon income received by the investor
    # Use coupon_rows count (not ac_idx) so it works even when the verified autocall
    # date doesn't exactly match a stored fecha_autocall_N due to format differences.
    if coupon_rows and cupon_period:
        rend_producto = len(coupon_rows) * cupon_period
    else:
        # Fallback: stored rendimiento_total from the portfolio spreadsheet
        rt_stored = _sf(product.get("rendimiento_total"))
        if rt_stored is not None:
            rend_producto = rt_stored * 100 if abs(rt_stored) <= 1 else rt_stored
        else:
            rend_producto = None

    # ── Narrative ──────────────────────────────────────────────────────────────
    narrative_paras = _narrative(
        name, event_type, actual_ac, fecha_inicio, fecha_vencto,
        ac_idx, cupon, barrera_pct, ganancia_max, factor_part, rend_producto
    )

    # ── Nivel Inicial / Final ──────────────────────────────────────────────────
    def _nivel(val_dict):
        parts = [f"{u}: {v:,.2f}" for u, v in val_dict.items() if v]
        if not parts:
            return "—"
        if len(parts) <= 2:
            return " / ".join(parts)
        mid = (len(parts) + 1) // 2
        return " / ".join(parts[:mid]) + "\n" + " / ".join(parts[mid:])

    # ── Subyacente description ─────────────────────────────────────────────────
    und_long = [_UND_LABELS.get(u, u) for u in unds]
    und_desc = ("Worst of: " + ", ".join(und_long)) if len(unds) > 1 else (und_long[0] if und_long else "—")
    subhdr_und = f"Subyacente: Opción Worst of sobre los ETF's de {', '.join(und_long)}" \
                 if len(unds) > 1 else f"Subyacente: {und_long[0] if und_long else '—'}"

    # ── Ganancia máxima ────────────────────────────────────────────────────────
    gmax = ganancia_max
    gmax_str = f"{gmax*100:.2f}%" if gmax and gmax <= 1 else (f"{gmax:.2f}%" if gmax else "—")

    # ── Características rows ───────────────────────────────────────────────────
    plazo_str = f"{int(plazo_meses)} meses" if plazo_meses else "—"
    moneda_str = "Dólares Americanos" if "USD" in moneda.upper() else moneda

    caract_rows = [
        ("Tipo de producto",         estrategia),
        ("Clasificación",            asset_class),
        ("Subyacente",               und_desc),
        ("Moneda",                   moneda_str),
        ("Plazo",                    plazo_str),
        ("Periodo sin Autocall",     periodo_nc),
        ("Observación Autocall",     obs_autocall),
        ("Observación Cupón",        obs_cupon),
        ("Cupón anual contingente",  _pct(cupon)),
        ("Ganancia Máxima",          gmax_str),
        ("Barrera",                  f"{barrera_pct:.0f}%"),
        ("Nivel Inicial",            _nivel(strikes)),
        ("Nivel Final",              _nivel(spots)),
        ("Código de Producto",       isin),
    ]
    caract_rows = [(k, v) for k, v in caract_rows if v and v not in ("—", "0%")]

    # ── Rendimiento subyacente for sub-header ─────────────────────────────────
    rend_p_str = f"{rend_producto:.3f}%" if rend_producto is not None else "—"
    worst_r_str = f"{worst_v:.2f}%" if worst_v is not None else "—"

    # ──────────────────────────────────────────────────────────────────────────
    # BUILD SLIDE
    # ──────────────────────────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(_W)
    prs.slide_height = Inches(_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    # 1. Header bar (always drawn)
    y0, hh = _YP["header"]
    _rect(slide, 0, y0, _W, hh, PRI)

    if logo_bytes:
        # Embed company logo on the left; height fills bar, width auto-scales
        slide.shapes.add_picture(
            BytesIO(logo_bytes),
            Inches(_ML),
            Inches(y0 + 0.04),
            width=None,
            height=Inches(hh - 0.08),
        )
    else:
        _txt(slide, _ML, y0 + 0.09, 2.5, hh - 0.18,
             company_name, size=9, bold=True, color=_WHITE)

    # Badge (event type label in top-right corner)
    bw = 1.0
    badge = _rect(slide, _W - bw - 0.20, y0 + 0.09, bw, hh - 0.18, _WHITE)
    btf = badge.text_frame
    btf.margin_left = btf.margin_right = Inches(0.05)
    btf.margin_top  = btf.margin_bottom = Inches(0.02)
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = event_type.upper()
    br.font.name = "Inter"; br.font.size = Pt(9)
    br.font.bold = True; br.font.color.rgb = PRI

    # 2. Title
    y0, hh = _YP["title"]
    _txt(slide, _ML, y0, _CW, hh,
         f"{event_type.upper()} – {name}",
         size=11, bold=True, color=SEC)

    # 3. DETALLE bar (secondary color — main section bar)
    y0, hh = _YP["det_bar"]
    _bar(slide, _ML, y0, _CW, hh, "DETALLE DEL PRODUCTO", SEC, size=9)

    # 4. Narrative
    y0, hh = _YP["narrative"]
    _multiline_txt(slide, _ML, y0, _CW, hh, narrative_paras, size=9)

    # 5. Split section bars
    y0, hh = _YP["sbars"]
    _bar(slide, _ML, y0, _LW, hh, "CARACTERÍSTICAS GENERALES", PRI, size=9)
    _bar(slide, _RX, y0, _RW, hh, "DESEMPEÑO DEL PEOR SUBYACENTE", PRI, size=9)

    # 6. Main two-column area
    y0, hh = _YP["main"]
    _build_caract_table(slide, _ML, y0, _LW, caract_rows, PRI)
    if not perf_df.empty:
        _embed(slide, _chart_all(
            perf_df,
            barrier_pct=barrier_cupon if cupon_annual else None,
            cupon_annual=cupon_annual if cupon_annual else None,
            accrual_color=secondary,
        ), _RX, y0, _RW, hh)

    # 7. Summary table
    y0, hh = _YP["sumtbl"]
    _build_summary_table(slide, _ML, y0, _CW, hh,
                         name, fecha_inicio, fecha_vencto, actual_ac,
                         rendimientos, worst_u, worst_v, rend_producto,
                         last_pmt, PRI)

    # 8. EVOLUCIÓN bar (secondary color — main section bar)
    y0, hh = _YP["evol_bar"]
    _bar(slide, _ML, y0, _CW, hh, "EVOLUCIÓN DEL PRODUCTO", SEC, size=9)

    # 9. Sub-header
    y0, hh = _YP["subhdr"]
    _txt(slide, _ML, y0, _CW * 0.64, hh, subhdr_und, size=9)
    _txt(slide, _ML + _CW * 0.66, y0,        _CW * 0.34, hh * 0.52,
         f"Rendimiento producto*: {rend_p_str}",
         size=9, bold=True, color=PRI, align=PP_ALIGN.RIGHT)
    _txt(slide, _ML + _CW * 0.66, y0 + hh * 0.52, _CW * 0.34, hh * 0.48,
         f"Rendimiento subyacente: {worst_r_str}",
         size=9, color=_DARK, align=PP_ALIGN.RIGHT)

    # 10. Bottom two-column
    y0, hh = _YP["bottom"]
    btm_lw = 3.00
    btm_rx = _ML + btm_lw + 0.15
    btm_rw = _CW - btm_lw - 0.15

    if coupon_rows:
        _build_coupon_table(slide, _ML, y0, btm_lw, coupon_rows, total_coupon_str, PRI)

    if not worst_series.empty:
        _embed(slide, _chart_worst(worst_series), btm_rx, y0, btm_rw, hh)

    # 11. Footnote / disclaimer
    y0, hh = _YP["footnote"]
    footnote_text = (
        disclaimer if disclaimer
        else "*Rentabilidad que se obtendría si el producto venciese en cada día observado."
    )
    _txt(slide, _ML, y0, _CW, hh,
         footnote_text, size=8, italic=not bool(disclaimer), color=_MUTED)

    # Serialize
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
